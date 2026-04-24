import sys
from pptx import Presentation

def inspect_ppt(file_path):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error loading pptx: {e}")
        return

    print("--- Slide Masters ---")
    for m_idx, master in enumerate(prs.slide_masters):
        print(f"\nMaster {m_idx}:")
        for s_idx, shape in enumerate(master.shapes):
            print(f"  Shape {s_idx}: type={shape.shape_type}, name='{shape.name}'")
            if hasattr(shape, 'image'):
                print(f"    Image: {shape.image.ext}")

        print(f"\n  Layouts in Master {m_idx}:")
        for l_idx, layout in enumerate(master.slide_layouts):
            print(f"    Layout {l_idx}: {layout.name}")
            for ls_idx, shape in enumerate(layout.shapes):
                print(f"      Shape {ls_idx}: type={shape.shape_type}, name='{shape.name}'")
                if shape.has_text_frame:
                    print(f"        Text: {repr(shape.text)}")
                if hasattr(shape, 'image'):
                    print(f"        Image: {shape.image.ext}")
                if hasattr(shape, 'click_action') and shape.click_action:
                    print(f"        Action: {shape.click_action.hyperlink.address if shape.click_action.hyperlink else 'None'}")

if __name__ == '__main__':
    inspect_ppt(sys.argv[1])
