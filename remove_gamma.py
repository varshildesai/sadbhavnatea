import sys
from pptx import Presentation

def remove_gamma_logo(input_path, output_path):
    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(f"Error loading pptx: {e}")
        return

    removed_count = 0
    
    # Check Masters and Layouts
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            # We must iterate backwards when deleting items
            for i in range(len(layout.shapes)-1, -1, -1):
                shape = layout.shapes[i]
                if hasattr(shape, 'click_action') and shape.click_action and shape.click_action.hyperlink:
                    if shape.click_action.hyperlink.address and 'gamma.app' in shape.click_action.hyperlink.address:
                        element = shape.element
                        element.getparent().remove(element)
                        removed_count += 1
                        print(f"Removed from layout {layout.name}")

    # Also check slides just in case
    for slide_idx, slide in enumerate(prs.slides):
        for i in range(len(slide.shapes)-1, -1, -1):
            shape = slide.shapes[i]
            if hasattr(shape, 'click_action') and shape.click_action and shape.click_action.hyperlink:
                if shape.click_action.hyperlink.address and 'gamma.app' in shape.click_action.hyperlink.address:
                    element = shape.element
                    element.getparent().remove(element)
                    removed_count += 1
                    print(f"Removed from slide {slide_idx}")

    prs.save(output_path)
    print(f"Done. Removed {removed_count} gamma logos. Saved to {output_path}")

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python remove_gamma.py input output")
    else:
        remove_gamma_logo(sys.argv[1], sys.argv[2])
